from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


BASE = Path(__file__).resolve().parent
OUT = BASE / "steady_study_proposal_v1.pptx"

W = Inches(13.333)
H = Inches(7.5)

COLORS = {
    "bg": RGBColor(246, 244, 240),
    "ink": RGBColor(18, 24, 33),
    "muted": RGBColor(95, 104, 117),
    "accent": RGBColor(255, 130, 22),
    "accent_dark": RGBColor(208, 95, 0),
    "panel": RGBColor(255, 255, 255),
    "panel_soft": RGBColor(252, 248, 242),
    "line": RGBColor(226, 220, 212),
    "emerald": RGBColor(14, 165, 90),
    "sky": RGBColor(38, 132, 255),
    "amber": RGBColor(245, 158, 11),
}


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_round_box(slide, x, y, w, h, fill, line=None, radius=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, name="Hiragino Sans"):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return tx


def add_bullets(slide, x, y, w, h, items, size=18, color=None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.font.name = "Hiragino Sans"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
    return tx


def add_tag(slide, x, y, w, text, fill, color=RGBColor(255, 255, 255)):
    box = add_round_box(slide, x, y, w, Inches(0.34), fill)
    box.line.color.rgb = fill
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Hiragino Sans"
    r.font.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = color
    return box


def header(slide, title, subtitle, tag_text=None):
    add_text(slide, Inches(0.55), Inches(0.38), Inches(8.9), Inches(0.5), title, size=24, bold=True)
    add_text(slide, Inches(0.55), Inches(0.86), Inches(10.2), Inches(0.32), subtitle, size=10, color=COLORS["muted"])
    if tag_text:
        add_tag(slide, Inches(11.2), Inches(0.36), Inches(1.55), tag_text, COLORS["accent"])


def footer(slide, text="Steady Study v1 / 2026"):
    add_text(slide, Inches(0.55), Inches(7.06), Inches(4.2), Inches(0.2), text, size=8, color=COLORS["muted"])


def section_card(slide, x, y, w, h, title, body, fill=COLORS["panel"], title_size=18, body_size=13):
    add_round_box(slide, x, y, w, h, fill, COLORS["line"])
    add_text(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.3), title, size=title_size, bold=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.48), w - Inches(0.36), h - Inches(0.58), body, size=body_size, color=COLORS["muted"])


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    layout = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    add_round_box(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.55), COLORS["panel"], COLORS["line"])
    add_tag(slide, Inches(0.7), Inches(0.74), Inches(2.15), "導入提案 / ONLINE塾向け", COLORS["accent"])
    add_text(slide, Inches(0.7), Inches(1.22), Inches(6.8), Inches(1.0), "Steady Study 導入提案 v1", size=30, bold=True)
    add_text(slide, Inches(0.7), Inches(2.15), Inches(6.8), Inches(0.65), "スマホ学習を 1 つに絞り、運用は半セルフで回す。", size=18, color=COLORS["muted"])
    add_text(slide, Inches(0.7), Inches(2.72), Inches(6.8), Inches(1.0), "今日やることは 1 つだけ。推奨コースは 1 つ。スペルチェックは既存のクイズに統合。", size=18, bold=True, color=COLORS["ink"])
    section_card(
        slide,
        Inches(8.0),
        Inches(1.1),
        Inches(4.3),
        Inches(4.8),
        "この提案で変えること",
        "・例文や画像を学習中に毎回生成しない\n・モバイルでは主推薦を 1 つだけ出す\n・導入相談は役割別の画面と分ける\n・管理者は導入チェックリストで開始できる",
        fill=COLORS["panel_soft"],
        title_size=16,
        body_size=15,
    )
    footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    header(slide, "塾が抱える課題", "公開ページと現場運用の両方から見た、導入時につまずきやすい点を整理しています。", "課題整理")
    section_card(slide, Inches(0.6), Inches(1.35), Inches(4.0), Inches(4.95), "教材準備が重い", "テキストは毎回 AI 生成ではなく、保存した例文を再利用したほうが安い。画像ヒントは学生導線では不要で、コストだけ上がりやすい。")
    section_card(slide, Inches(4.67), Inches(1.35), Inches(4.0), Inches(4.95), "生徒が迷いやすい", "画面に選択肢が多いと、スマホでは何から始めるか分かりにくい。学習の入口は 1 つ、推奨コースも 1 つに絞る必要がある。")
    section_card(slide, Inches(8.74), Inches(1.35), Inches(4.0), Inches(4.95), "導入運用の手間", "オンライン中小塾では、最初の設定と割当が重いと止まる。公開相談、役割確認、開始手順を分けて半セルフ化するのが前提。")
    footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    header(slide, "Steady Study の解き方", "公開導線は 4 ステップに固定し、導入相談から開始までの迷いを減らします。", "4 steps")
    step_titles = ["1. 対象塾か判断", "2. 役割別の見え方確認", "3. 導入相談を送る", "4. 招待または手動 provisioning"]
    step_bodies = [
        "オンライン中心か、教材配信したいかを先に確認します。",
        "生徒・講師・管理者の画面を役割別に見せます。",
        "授業形態、開始時期、想定人数だけを送ってもらいます。",
        "自動発行はせず、案内後に個別 provisioning で開始します。",
    ]
    xs = [0.55, 3.7, 6.85, 10.0]
    fills = [COLORS["accent"], COLORS["sky"], COLORS["emerald"], COLORS["amber"]]
    for i in range(4):
        add_round_box(slide, Inches(xs[i]), Inches(1.55), Inches(2.7), Inches(3.95), COLORS["panel"], COLORS["line"])
        add_tag(slide, Inches(xs[i] + 0.18), Inches(1.77), Inches(1.82), step_titles[i], fills[i])
        add_text(slide, Inches(xs[i] + 0.18), Inches(2.22), Inches(2.34), Inches(0.55), step_bodies[i], size=16, bold=True)
    add_round_box(slide, Inches(0.7), Inches(5.82), Inches(11.95), Inches(0.72), COLORS["panel_soft"], COLORS["line"])
    add_text(slide, Inches(0.95), Inches(6.03), Inches(11.45), Inches(0.22), "営業トークは「半セルフ導入」。公開説明は最小限にして、相談と個別案内を主導線にします。", size=14, color=COLORS["muted"])
    footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    header(slide, "画面の約束", "生徒・講師・管理者それぞれの画面で、見せる情報を最小限に整理します。", "UI")
    section_card(slide, Inches(0.55), Inches(1.4), Inches(4.0), Inches(4.9), "生徒画面", "・今日やることは 1 つだけ\n・主推薦コースは 1 つだけ\n・学習カードの例文は保存済みを表示\n・スペルチェックは全文入力→ヒントの順", fill=COLORS["panel"])
    section_card(slide, Inches(4.67), Inches(1.4), Inches(4.0), Inches(4.9), "講師画面", "・公式単語帳と My単語帳を切り替え\n・学習状況、停滞リスク、配布を一画面で把握\n・教材に対して例文準備を明示実行", fill=COLORS["panel"])
    section_card(slide, Inches(8.79), Inches(1.4), Inches(3.95), Inches(4.9), "管理者画面", "・導入相談キューを確認\n・AI 利用は集計を正史に統一\n・組織名、cohort、講師割当、教材配布、初回ミッションを確認", fill=COLORS["panel"])
    footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    header(slide, "AIコスト方針", "学習中に AI を呼ばない設計へ寄せ、保存再利用を前提にしています。", "cost")
    section_card(slide, Inches(0.6), Inches(1.45), Inches(6.0), Inches(4.9), "採用方針", "・例文は教材準備時にだけ作る\n・学習中は `words.example_sentence` / `example_meaning` を再利用\n・学生導線の画像生成は v1 では外す\n・今月の利用表示は固定モックではなく実集計に寄せる", fill=COLORS["panel_soft"], body_size=16)
    section_card(slide, Inches(6.8), Inches(1.45), Inches(5.0), Inches(4.9), "運用上の意味", "・例文の保存再利用は、毎回生成するより明確に安い\n・画像生成は特に高コストなので学生フローから外す\n・必要なときだけ管理者が例文準備を実行する", fill=COLORS["panel"])
    footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(layout)
    set_bg(slide, COLORS["bg"])
    header(slide, "導入ステップ", "最初の 1 回で迷わないために、開始時の確認事項を固定化しています。", "start")
    section_card(slide, Inches(0.6), Inches(1.42), Inches(12.0), Inches(3.85), "初回チェックリスト", "1. 組織名を確認する\n2. cohort を作成する\n3. 講師を割り当てる\n4. 教材を配布する\n5. 初回ミッションを設定する", fill=COLORS["panel"])
    add_round_box(slide, Inches(0.6), Inches(5.55), Inches(12.0), Inches(0.95), COLORS["accent"])
    add_text(slide, Inches(0.9), Inches(5.83), Inches(11.4), Inches(0.3), "次のアクション: 1) 公開ページで役割を確認 2) 導入相談送信 3) 管理者が開始設定", size=16, bold=True, color=RGBColor(255, 255, 255))
    footer(slide)

    prs.save(OUT)


if __name__ == "__main__":
    build()
